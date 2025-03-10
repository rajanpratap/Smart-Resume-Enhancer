# resumeapp/resume_analysis.py

import time, requests
import os
from bs4 import BeautifulSoup
from langchain_groq import ChatGroq
from langchain.prompts import PromptTemplate
from langgraph.graph import StateGraph, END
from docx import Document
from pptx import Presentation
import fitz  # PyMuPDF for PDFs
from difflib import ndiff
from dotenv import load_dotenv

load_dotenv()
temp_dir = "temp_files"
os.makedirs(temp_dir, exist_ok=True)

GROQ_API_KEY = os.getenv("GROQ_API_KEY")
print(GROQ_API_KEY)

# Set up LLM
llm = ChatGroq(model_name="gemma2-9b-it", temperature=0.7, api_key=GROQ_API_KEY)

def rebuild_resume_with_format(updated_text, formatting, file_type):
    file_path = os.path.join(temp_dir, f"updated_resume.{file_type}")
    
    if file_type == "docx":
        doc = Document()
        for entry in formatting:
            para = doc.add_paragraph()
            run = para.add_run(entry["text"])
            if entry.get("bold"):
                run.bold = True
            if entry.get("italic"):
                run.italic = True
        doc.save(file_path)
    elif file_type == "pptx":
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        textbox = slide.shapes.add_textbox(100, 100, 500, 300)
        text_frame = textbox.text_frame
        text_frame.text = updated_text
        prs.save(file_path)
    elif file_type == "pdf":
        doc = fitz.open()
        page = doc.new_page()
        page.insert_text((50, 50), updated_text)
        doc.save(file_path)
    
    return file_path

# Define state structure (for clarity; not used directly in Django)
class ResumeUpdateState:
    def __init__(self, resume: str, job_description: str, extracted_requirements="", resume_gaps="", updated_resume=""):
        self.resume = resume
        self.job_description = job_description
        self.extracted_requirements = extracted_requirements
        self.resume_gaps = resume_gaps
        self.updated_resume = updated_resume

def extract_text_and_format_from_resume(file, file_type):
    file_path = os.path.join(temp_dir, file.name)
    
    with open(file_path, "wb") as f:
        f.write(file.read())  # In Django, use file.read()
    
    extracted_text = ""
    formatting = []
    
    if file_type == "docx":
        doc = Document(file_path)
        for para in doc.paragraphs:
            extracted_text += para.text + "\n"
            formatting.append({
                "text": para.text,
                "bold": any(run.bold for run in para.runs),
                "italic": any(run.italic for run in para.runs),
                "font_size": para.runs[0].font.size if para.runs else None
            })
    elif file_type == "pptx":
        prs = Presentation(file_path)
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    extracted_text += shape.text + "\n"
                    formatting.append({"text": shape.text, "font_size": None})
    elif file_type == "pdf":
        doc = fitz.open(file_path)
        for page in doc:
            for text_block in page.get_text("dict")["blocks"]:
                if "lines" in text_block:
                    for line in text_block["lines"]:
                        for span in line["spans"]:
                            extracted_text += span["text"] + " "
                            formatting.append({"text": span["text"], "font_size": span["size"]})
            extracted_text += "\n"
    else:
        return "", []
    
    safe_remove(file_path)
    return extracted_text.strip(), formatting

def safe_remove(file_path):
    for _ in range(5):
        try:
            os.remove(file_path)
            break
        except PermissionError:
            time.sleep(0.5)

def fetch_job_description_from_llm(url):
    try:
        headers = {"User-Agent": "Mozilla/5.0"}
        response = requests.get(url, headers=headers)
        response.raise_for_status()

        soup = BeautifulSoup(response.text, "html.parser")
        page_text = soup.get_text(separator=" ", strip=True)

        job_description = llm.invoke(f"""
        Extract the job description from the following webpage text. Focus only on the responsibilities, qualifications, and skills mentioned for the role. 
        Ignore navigation menus, company descriptions, and unrelated content. Show each skill or requirement on a new line.

        Webpage Content:
        {page_text}

        Extracted Job Description:
        """)
        return job_description
    except Exception as e:
        print(f"Error fetching job description: {str(e)}")
        return None

def highlight_differences(original, updated):
    diff = list(ndiff(original.split(), updated.split()))
    highlighted = " ".join([
        f'**{word[2:]}**' if word.startswith('+ ') else word[2:]
        for word in diff if not word.startswith('- ')
    ])
    return highlighted

def extract_job_requirements(state):
    extracted_text = llm.invoke(extract_prompt.format(job_description=state["job_description"]))
    return {
        "resume": state["resume"],
        "job_description": state["job_description"],
        "extracted_requirements": extracted_text,
        "formatting": state["formatting"],
        "file_type": state["file_type"]
    }

def analyze_resume_gaps(state):
    gap_analysis = llm.invoke(analyze_prompt.format(resume=state["resume"], extracted_requirements=state["extracted_requirements"]))
    return {
        "resume": state["resume"],
        "job_description": state["job_description"],
        "extracted_requirements": state["extracted_requirements"],
        "resume_gaps": gap_analysis.content,
        "formatting": state["formatting"],
        "file_type": state["file_type"]
    }

def update_resume(state):
    updated_text = llm.invoke(f"""
        Modify the resume below to align with the job requirements while keeping it truthful and professional.
        Highlight missing skills and reframe experience to better match the role.
        
        Resume:
        {state["resume"]}
        
        Identified Gaps:
        {state["resume_gaps"]}
        
        Where possible, subtly incorporate missing but relevant skills into existing experience sections.
        Return only the modified text against original resume
    """)
    highlighted_text = highlight_differences(state["resume"], updated_text.content)
    return {
        "resume": state["resume"],
        "job_description": state["job_description"],
        "extracted_requirements": state["extracted_requirements"],
        "resume_gaps": state["resume_gaps"],
        "updated_resume": highlighted_text
    }

# Prompt templates for LLM
extract_prompt = PromptTemplate(
    input_variables=["job_description"],
    template="""
    Extract key skills, responsibilities, and qualifications from the following job description:
    {job_description}
    """
)

analyze_prompt = PromptTemplate(
    input_variables=["resume", "extracted_requirements"],
    template="""
    Compare the given resume with the extracted job requirements and identify missing skills, experience gaps, and areas of improvement.
    Resume:
    {resume}
    Job Requirements:
    {extracted_requirements}
    """
)

# Build the state graph for the resume updater
graph = StateGraph(dict)
graph.add_node("extract", extract_job_requirements)
graph.add_node("analyze", analyze_resume_gaps)
graph.add_node("update", update_resume)

graph.add_edge("extract", "analyze")
graph.add_edge("analyze", "update")
graph.add_edge("update", END)

graph.set_entry_point("extract")
graph = graph.compile()

def run_resume_analysis(uploaded_file, job_url):
    # Determine file type from file name
    file_type = uploaded_file.name.split(".")[-1].lower()
    resume_text, formatting = extract_text_and_format_from_resume(uploaded_file, file_type)
    if not resume_text:
        return {"error": "Unsupported file type or extraction error."}
    
    job_description_obj = fetch_job_description_from_llm(job_url)
    if not job_description_obj:
        return {"error": "Error fetching job description."}
    
    job_description = job_description_obj.content if hasattr(job_description_obj, "content") else str(job_description_obj)
    
    state = {
        "resume": resume_text,
        "job_description": job_description,
        "formatting": formatting,
        "file_type": file_type
    }
    output = graph.invoke(state)
    return output
